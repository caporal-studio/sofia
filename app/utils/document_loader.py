import os
import fitz
import pandas as pd
from docx import Document
from pptx import Presentation
from bs4 import BeautifulSoup
import json
import xml.etree.ElementTree as ET
from app.utils.app_config import get_app_info

info = get_app_info()
modelo_configurado = info.get("openai_model") or info.get("ollama_model", "gpt-4o")

# Variáveis de importação lazy load
DocumentConverter = None
HybridChunker = None
enc = None

docling_available = False
imports_checked = False

def setup_imports():
    global DocumentConverter, HybridChunker, enc, docling_available, imports_checked
    if imports_checked:
        return
    imports_checked = True

    try:
        from docling.document_converter import DocumentConverter as DC
        from docling.chunking import HybridChunker as HC
        DocumentConverter = DC
        HybridChunker = HC
        docling_available = True
        print("✅ Docling importado com sucesso.")
    except ImportError:
        docling_available = False
    except Exception as e:
        print(f"Docling opcional indisponivel: {e}")
        docling_available = False

    try:
        import tiktoken
        try:
            enc = tiktoken.encoding_for_model(modelo_configurado)
        except KeyError:
            enc = tiktoken.get_encoding("cl100k_base")
        print("✅ Tiktoken importado com sucesso.")
    except ImportError:
        enc = None
    except Exception as e:
        print(f"Tiktoken indisponivel: {e}")
        enc = None

# ======= Chunking Inteligente (docling + tiktoken) ======= #
def chunk_text_with_docling(path, max_tokens=750, overlap=100):
    setup_imports()
    if not docling_available or DocumentConverter is None or HybridChunker is None:
        return None

    try:
        converter = DocumentConverter()
        result = converter.convert(source=path)
        doc = result.document
        try:
            from transformers import AutoTokenizer
        except ImportError:
            print("Docling instalado, mas transformers nao esta disponivel. Usando fallback por tokens.")
            return None
        tokenizer = AutoTokenizer.from_pretrained("sentence-transformers/all-MiniLM-L6-v2", model_max_length=512, truncation=True)
        chunker = HybridChunker(tokenizer=tokenizer, max_tokens=max_tokens, overlap=overlap)
        chunks = [chunker.contextualize(chunk).strip() for chunk in chunker.chunk(dl_doc=doc)]
        print(f"✅ Docling processou com sucesso: {path}")
        return [chunk for chunk in chunks if chunk]
    except Exception as e:
        print(f"⚠️ Docling falhou em {path}: {e}")
        return None

def chunk_text_with_tiktoken(text, max_tokens=750, overlap=100):
    global enc
    final_chunks = []

    if enc is None:
        try:
            import tiktoken
            enc = tiktoken.encoding_for_model(modelo_configurado)
            print("✅ Tiktoken importado sob demanda.")
        except Exception as e:
            print(f"⚠️ Erro ao importar tiktoken: {e}")
            enc = None

    if enc:
        tokens = enc.encode(text)
        for i in range(0, len(tokens), max_tokens - overlap):
            sub_tokens = tokens[i:i + max_tokens]
            sub_chunk = enc.decode(sub_tokens)
            final_chunks.append(sub_chunk)
            print("✅ Controle de tokens do Tiktoken ativado.")
    else:
        print("⚠️ Usando fallback por parágrafos, Tiktoken indisponível.")
        words = text.split()
        for i in range(0, len(words), max_tokens - overlap):
            sub_chunk = " ".join(words[i:i + max_tokens])
            final_chunks.append(sub_chunk)

    return [c.strip() for c in final_chunks if c.strip()]

# ======= Funções de extração por tipo de arquivo ======= #
def extract_text_from_pdf(path):
    doc = fitz.open(path)
    return "\n".join(page.get_text() for page in doc)

def extract_text_from_docx(path):
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs)

def extract_text_from_xlsx(path):
    try:
        dfs = pd.read_excel(path, sheet_name=None)
        full_text = ""
        for sheet_name, df in dfs.items():
            df.columns = df.columns.map(str)
            full_text += f"📄 Planilha: {os.path.basename(path)} | Aba: {sheet_name}\n"
            full_text += df.to_markdown(index=False) + "\n\n"
        return full_text
    except Exception as e:
        print(f"Erro ao processar {path}: {e}")
        return ""

def extract_text_from_pptx(path):
    prs = Presentation(path)
    return "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))

def extract_text_from_txt(path):
    with open(path, "r", encoding="utf-8") as f:
        return f.read()

def extract_text_from_json(path):
    with open(path, "r", encoding="utf-8") as f:
        data = json.load(f)
    return json.dumps(data, indent=2, ensure_ascii=False)

def extract_text_from_csv(path):
    df = pd.read_csv(path)
    return df.to_string(index=False)

def extract_text_from_xml(path):
    tree = ET.parse(path)
    root = tree.getroot()
    return ET.tostring(root, encoding='unicode')

def extract_text_from_md(path):
    with open(path, "r", encoding="utf-8") as f:
        return f.read()

def extract_text_from_html(path):
    with open(path, "r", encoding="utf-8") as f:
        soup = BeautifulSoup(f.read(), "html.parser")
    return soup.get_text()

def extract_text_from_image(path):
    from PIL import Image
    import pytesseract

    image = Image.open(path)
    return pytesseract.image_to_string(image)

# ======= Carregamento e chunking dos documentos ======= #
def load_documents(folder_path="documentacao"):
    documents = []
    supported_extensions = {
        ".pdf": extract_text_from_pdf,
        ".docx": extract_text_from_docx,
        ".xlsx": extract_text_from_xlsx,
        ".pptx": extract_text_from_pptx,
        ".txt": extract_text_from_txt,
        ".json": extract_text_from_json,
        ".csv": extract_text_from_csv,
        ".xml": extract_text_from_xml,
        ".md": extract_text_from_md,
        ".html": extract_text_from_html,
        ".htm": extract_text_from_html,
        ".png": extract_text_from_image,
        ".jpg": extract_text_from_image,
        ".jpeg": extract_text_from_image,
    }

    for root, _, files in os.walk(folder_path):
        for file in files:
            ext = os.path.splitext(file)[1].lower()
            extractor = supported_extensions.get(ext)
            if not extractor:
                continue

            file_path = os.path.join(root, file)
            try:
                file_rel_path = os.path.relpath(file_path, folder_path).replace("\\", "/")

                chunks = chunk_text_with_docling(file_path)
                if chunks is None:
                    content = extractor(file_path)
                    chunks = chunk_text_with_tiktoken(content)

                for i, chunk in enumerate(chunks):
                    documents.append({
                        "source": f"{file_rel_path} [parte {i+1}]",
                        "content": chunk
                    })

            except Exception as e:
                print(f"Erro ao processar {file_path}: {e}")

    return documents

def load_tabular_data(folder_path="documentacao"):
    tabular_data = {}

    for root, _, files in os.walk(folder_path):
        for file in files:
            ext = os.path.splitext(file)[1].lower()
            full_path = os.path.join(root, file)

            try:
                if ext == ".xlsx":
                    df = pd.read_excel(full_path)
                elif ext == ".csv":
                    df = pd.read_csv(full_path)
                elif ext == ".json":
                    df = pd.read_json(full_path)
                else:
                    continue

                relative_name = os.path.relpath(full_path, folder_path).replace("\\", "/")
                tabular_data[relative_name] = df

            except Exception as e:
                print(f"Erro ao carregar {file}: {e}")

    return tabular_data
